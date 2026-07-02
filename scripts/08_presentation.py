"""
Generate the project presentation as a PowerPoint file.

Sections:
  1. Project Goal
  2. The Base Paper (Li et al. 2022)
  3. Our Project
  4. Approach & Datasets
  5. Features of Interest
  6. Project Plan (high-level)
  7. Phase Details
  8. Results
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE_DIR = Path(__file__).resolve().parent.parent
IMG_DIR = BASE_DIR / "data" / "processed" / "analysis"
OUT_PATH = BASE_DIR / "presentation.pptx"

# ── Colors ──────────────────────────────────────────────────────────────────
DARK = RGBColor(0x2C, 0x3E, 0x50)
ACCENT = RGBColor(0x27, 0xAE, 0x60)
LIGHT_BG = RGBColor(0xEC, 0xF0, 0xF1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
RED = RGBColor(0xC0, 0x39, 0x2B)
BLUE = RGBColor(0x29, 0x80, 0xB9)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _fmt_para(p, text, font_size, bold, color, font_name, alignment):
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    _fmt_para(tf.paragraphs[0], lines[0], font_size, bold, color, font_name, alignment)
    for line in lines[1:]:
        p = tf.add_paragraph()
        _fmt_para(p, line, font_size, bold, color, font_name, alignment)
    return tf


def add_bullet_slide(slide, left, top, width, height, items, font_size=16,
                     color=DARK, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        lines = item.split("\n")
        for j, line in enumerate(lines):
            if i == 0 and j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = "Calibri"
            p.space_after = spacing if j == len(lines) - 1 else Pt(0)
            p.level = 0
    return tf


def add_section_header(slide, number, title):
    set_slide_bg(slide, DARK)
    # Number circle
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(1), Inches(1),
                 str(number), font_size=48, bold=True, color=ACCENT,
                 alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, Inches(2.0), Inches(1.5), Inches(7), Inches(1.2),
                 title, font_size=36, bold=True, color=WHITE)


def add_content_title(slide, title):
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
                 title, font_size=28, bold=True, color=DARK)
    # Accent underline
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0.6), Inches(0.95), Inches(2), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    blank = prs.slide_layouts[6]  # blank layout

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # TITLE SLIDE
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, DARK)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.2),
                 "Smart Walkability Perception",
                 font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(8.4), Inches(0.8),
                 "Measuring How Walkable Tel Aviv's Streets Feel\nUsing Municipal Open Data",
                 font_size=22, color=ACCENT)
    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.5),
                 "Doron Peleg  |  2026", font_size=14, color=GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 1. PROJECT GOAL
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 1, "Project Goal")

    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "What Makes a Street Feel Walkable?")
    add_bullet_slide(slide, Inches(0.6), Inches(1.3), Inches(8.8), Inches(3.5), [
        "Some streets invite you to walk. Others make you want to drive.",
        "This \"feeling\" is called Visual Walkability Perception (VWP) —\n"
        "it combines safety, comfort, liveliness, and accessibility into\n"
        "a single impression of how pleasant a street is for pedestrians.",
        "Traditional measurement requires expensive surveys or\n"
        "street-level imagery analysis. We asked a different question:",
    ], font_size=16)
    add_text_box(slide, Inches(1.0), Inches(4.0), Inches(8), Inches(0.7),
                 "Can a city's own operational data — business records, complaints,\n"
                 "demographics — tell us how walkable its streets are?",
                 font_size=18, bold=True, color=ACCENT)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 2. THE BASE PAPER
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 2, "The Base Paper")

    # Paper overview
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Li et al. (2022) — Measuring Visual Walkability")
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.5),
                 "The Problem", font_size=20, bold=True, color=BLUE)
    add_text_box(slide, Inches(0.6), Inches(1.7), Inches(8.8), Inches(0.7),
                 "How walkable a street looks matters for urban planning, health,\n"
                 "and property values — but measuring it at city scale is hard.\n"
                 "Surveys are expensive. Audits are subjective.",
                 font_size=15, color=DARK)

    add_text_box(slide, Inches(0.6), Inches(2.6), Inches(8.8), Inches(0.5),
                 "Their Method", font_size=20, bold=True, color=BLUE)
    add_bullet_slide(slide, Inches(0.6), Inches(3.1), Inches(8.8), Inches(1.5), [
        "Collected 2,642 panoramic street images from 7 world cities",
        "30 people rated them in VR headsets across 6 walkability dimensions",
        "Used AI (deep learning) to identify what makes streets look walkable",
        "Found that vegetation, sidewalks, and traffic signs matter most",
    ], font_size=14, spacing=Pt(3))

    add_text_box(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(0.5),
                 "Key Result:  19 visual features explain 41–59% of walkability perception",
                 font_size=15, bold=True, color=ACCENT)

    # Paper VWP categories
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Six Dimensions of Walkability")
    cats = [
        ("Walkability", "Overall impression — does this street support walking?"),
        ("Feasibility", "Is there a reason to walk here? Shops, services, mixed land use"),
        ("Accessibility", "Can you walk here? No barriers, dead ends, or missing sidewalks"),
        ("Safety", "Does it feel safe from crime and traffic?"),
        ("Comfort", "Is it pleasant? Wide sidewalks, shade, low noise"),
        ("Pleasurability", "Is it enjoyable? Greenery, lively atmosphere, interesting views"),
    ]
    for i, (cat, desc) in enumerate(cats):
        y = 1.2 + i * 0.65
        add_text_box(slide, Inches(0.8), Inches(y), Inches(2.2), Inches(0.5),
                     cat, font_size=16, bold=True, color=ACCENT)
        add_text_box(slide, Inches(3.0), Inches(y), Inches(6.5), Inches(0.5),
                     desc, font_size=14, color=DARK)

    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.4),
                 "Based on Alfonzo's Hierarchy of Walking Needs (2005)",
                 font_size=11, color=GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 3. OUR PROJECT
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 3, "Our Project")

    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "A Different Lens on Walkability")

    add_text_box(slide, Inches(0.6), Inches(1.3), Inches(4.2), Inches(0.4),
                 "The paper used:", font_size=16, bold=True, color=GRAY)
    add_bullet_slide(slide, Inches(0.6), Inches(1.7), Inches(4.2), Inches(1.2), [
        "Street-level photos",
        "AI to count pixels of trees, cars, signs",
        "Human raters in VR headsets",
    ], font_size=14, color=GRAY)

    add_text_box(slide, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.4),
                 "We use instead:", font_size=16, bold=True, color=ACCENT)
    add_bullet_slide(slide, Inches(5.2), Inches(1.7), Inches(4.5), Inches(1.2), [
        "Tel Aviv municipal open data",
        "Business registries, complaint logs, demographics",
        "Traffic speed as walkability proxy",
    ], font_size=14, color=DARK)

    add_text_box(slide, Inches(0.6), Inches(3.2), Inches(8.8), Inches(0.5),
                 "Key Insight", font_size=20, bold=True, color=DARK)
    add_text_box(slide, Inches(0.6), Inches(3.7), Inches(8.8), Inches(1.5),
                 "A camera sees trees, sidewalks, and cars.\n"
                 "Municipal data sees what the camera can't:\n"
                 "•  Are businesses thriving or closing?\n"
                 "•  Are residents complaining about noise and broken sidewalks?\n"
                 "•  Do women feel safe enough to move into the neighborhood?\n"
                 "•  Are there cultural events that bring people to the street?",
                 font_size=15, color=DARK)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 4. APPROACH & DATASETS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 4, "Approach & Datasets")

    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Tel Aviv Open Data")
    datasets = [
        ("Traffic Speed Segments", "24 months of hourly speeds on 2,070 road segments", "ArcGIS REST API"),
        ("Business Registry", "35,172 businesses with category, licensing, neighborhood", "Azure Blob CSV"),
        ("106 Hotline Complaints", "286,657 municipal complaints (noise, infrastructure...)", "Azure Blob CSV"),
        ("Closed Streets", "2,006 street closures: events, construction, light rail", "Azure Blob CSV"),
        ("Population & Migration", "Age, gender, neighborhood breakdown (2010–2022)", "CBS/Municipal CSV"),
        ("Dwelling Units", "Residential vs. commercial unit counts by neighborhood", "Municipal CSV"),
        ("Socio-Economic Index", "18 indicators per neighborhood (income, cars, education)", "Municipal Excel"),
        ("Construction Starts", "Building starts/completions by district (2020–2022)", "CBS CSV"),
    ]
    for i, (name, desc, source) in enumerate(datasets):
        y = 1.15 + i * 0.5
        add_text_box(slide, Inches(0.6), Inches(y), Inches(2.8), Inches(0.45),
                     name, font_size=12, bold=True, color=DARK)
        add_text_box(slide, Inches(3.4), Inches(y), Inches(4.6), Inches(0.45),
                     desc, font_size=11, color=DARK)
        add_text_box(slide, Inches(8.0), Inches(y), Inches(1.8), Inches(0.45),
                     source, font_size=9, color=GRAY)

    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.4),
                 "Total: 35 raw files  |  255 MB  |  1.3M+ records  |  All freely available",
                 font_size=13, bold=True, color=ACCENT)

    # Approach diagram slide
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Methodology")
    steps = [
        ("Download\n& Clean", "8 data sources\n35 files\nUTF-8 normalization"),
        ("Spatial\nDatabase", "GeoPackage\n16 layers\n71 neighborhoods"),
        ("Feature\nEngineering", "27 features\nper segment\n5 join strategies"),
        ("Modeling\n& PCA", "Linear, RF, NN\nPermutation imp.\nComposite score"),
        ("Validation\n& Maps", "Paper comparison\nInteractive map\nRecommendations"),
    ]
    for i, (title, desc) in enumerate(steps):
        x = 0.3 + i * 1.95
        is_highlight = (i == 3)
        fg = WHITE if is_highlight else DARK
        fg2 = WHITE if is_highlight else GRAY
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x), Inches(1.6), Inches(1.7), Inches(2.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT if is_highlight else LIGHT_BG
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        for j, line in enumerate(title.split("\n")):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            _fmt_para(p, line, 14, True, fg, "Calibri", PP_ALIGN.CENTER)
        for j, line in enumerate(desc.split("\n")):
            p = tf.add_paragraph()
            _fmt_para(p, line, 10, False, fg2, "Calibri", PP_ALIGN.CENTER)
            if j == 0:
                p.space_before = Pt(10)
        if i < len(steps) - 1:
            add_text_box(slide, Inches(x + 1.72), Inches(2.4), Inches(0.3), Inches(0.4),
                         "→", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(4.2), Inches(8.8), Inches(0.8),
                 "Unit of analysis: directed traffic segment (~100m of street)\n"
                 "Outcome variable: average traffic speed (proxy for walkability)",
                 font_size=14, color=DARK)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 5. FEATURES
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 5, "Features of Interest")

    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "27 Features Across 8 Categories")
    feature_groups = [
        ("Traffic (5)", "avg speed, peak-hour drop, speed variance, directional differential, evening recovery", BLUE),
        ("Business (5)", "restaurant count, total businesses, diversity index, licensed ratio, evening business count", ACCENT),
        ("Street Events (5)", "event frequency, event diversity, light rail closures, parade routes, infrastructure closures", RGBColor(0x8E, 0x44, 0xAD)),
        ("Complaints (2)", "noise complaints, infrastructure complaints", RED),
        ("Construction (2)", "building starts, completion rate", RGBColor(0xD3, 0x54, 0x00)),
        ("Demographics (3)", "female population growth, young adult share (20-34), elderly share (65+)", RGBColor(0x00, 0x97, 0xA7)),
        ("Housing (2)", "residential density, commercial-to-residential ratio", RGBColor(0x79, 0x55, 0x48)),
        ("Socio-Economic (3)", "SES cluster, cars per 100 residents, avg monthly income", GRAY),
    ]
    for i, (cat, feats, color) in enumerate(feature_groups):
        row = i // 2
        col = i % 2
        x = 0.5 + col * 4.8
        y = 1.2 + row * 0.95
        add_text_box(slide, Inches(x), Inches(y), Inches(1.8), Inches(0.4),
                     cat, font_size=12, bold=True, color=color)
        add_text_box(slide, Inches(x + 1.8), Inches(y), Inches(3.0), Inches(0.8),
                     feats, font_size=10, color=DARK)

    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.4),
                 "Coverage: 82–100% across all features  |  Joined via spatial + code-based strategies",
                 font_size=12, color=GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 6. PROJECT PLAN
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 6, "Project Plan")

    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Four Phases, Eight Weeks")
    phases = [
        ("Phase 1", "Data Collection\n& Pipeline", "Weeks 1–2",
         "Download, clean, and store\nall datasets in a spatial database"),
        ("Phase 2", "Feature\nEngineering", "Weeks 3–4",
         "Compute 27 features per\nstreet segment from raw data"),
        ("Phase 3", "Analysis &\nModeling", "Weeks 5–6",
         "Regression, PCA, walkability\nscoring and mapping"),
        ("Phase 4", "Validation &\nInterpretation", "Weeks 7–8",
         "Compare with paper, generate\nrecommendations and deliverables"),
    ]
    for i, (phase, title, timeline, desc) in enumerate(phases):
        x = 0.3 + i * 2.4
        # Phase box
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x), Inches(1.4), Inches(2.1), Inches(3.0))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_BG
        shape.line.fill.background()
        # Phase number
        add_text_box(slide, Inches(x + 0.1), Inches(1.5), Inches(1.9), Inches(0.4),
                     phase, font_size=12, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(1.9), Inches(1.9), Inches(0.6),
                     title, font_size=15, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(2.5), Inches(1.9), Inches(0.3),
                     timeline, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(x + 0.1), Inches(2.9), Inches(1.9), Inches(1.2),
                     desc, font_size=10, color=DARK, alignment=PP_ALIGN.CENTER)
        # Arrow
        if i < 3:
            add_text_box(slide, Inches(x + 2.12), Inches(2.5), Inches(0.3), Inches(0.4),
                         "→", font_size=20, color=ACCENT, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.4),
                 "Tech stack: Python, GeoPandas, statsmodels, scikit-learn, Folium  |  "
                 "Storage: GeoPackage (243 MB)",
                 font_size=11, color=GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 7. PHASE DETAILS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 7, "Phase Details")

    # Phase 1 detail
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Phase 1: Data Collection & Pipeline")
    add_bullet_slide(slide, Inches(0.6), Inches(1.2), Inches(5.0), Inches(4.0), [
        "Downloaded 24 monthly traffic layers from ArcGIS REST API\n"
        "with pagination (2,000 records/request)",
        "Downloaded 10 CSV datasets, handling 3 different encodings\n"
        "(UTF-16LE, Windows-1255, UTF-8) and 4 different delimiters",
        "Normalized all column names: Hebrew → English snake_case",
        "Built GeoPackage spatial database with 16 layers",
        "Downloaded 71 neighborhood boundary polygons\n"
        "for spatial joins in later phases",
    ], font_size=13, spacing=Pt(4))
    add_text_box(slide, Inches(5.8), Inches(1.3), Inches(4.0), Inches(0.4),
                 "Output:", font_size=14, bold=True, color=ACCENT)
    add_bullet_slide(slide, Inches(5.8), Inches(1.7), Inches(4.0), Inches(2.5), [
        "35 raw files (255 MB)",
        "50,947 traffic segments (24 months)",
        "1.3M+ records across all datasets",
        "smarttraffic.gpkg (243 MB, 16 layers)",
    ], font_size=12, color=DARK)
    img = IMG_DIR.parent / "traffic_segments_map.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(6.2), Inches(3.2), Inches(2.8))

    # Phase 2 detail
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Phase 2: Feature Engineering")
    add_bullet_slide(slide, Inches(0.6), Inches(1.2), Inches(5.0), Inches(4.0), [
        "Spatial join: segment centroids → neighborhood polygons\n"
        "(99.6% matched, using projected CRS for accuracy)",
        "5 traffic features computed directly from hourly speed columns",
        "5 business features via neighborhood_code join\n"
        "(Shannon entropy for diversity, license ratio, evening count)",
        "5 street event features via street_code join\n"
        "(event frequency, light rail, parade routes)",
        "Complaint, construction, demographic, housing, SES features\n"
        "aggregated at neighborhood level with name normalization",
    ], font_size=13, spacing=Pt(4))
    add_text_box(slide, Inches(5.8), Inches(1.3), Inches(4.0), Inches(0.4),
                 "Output:", font_size=14, bold=True, color=ACCENT)
    add_bullet_slide(slide, Inches(5.8), Inches(1.7), Inches(4.0), Inches(2.0), [
        "2,070 segments × 27 features",
        "Coverage: 82–100% per feature",
        "'features' layer in GeoPackage",
    ], font_size=12, color=DARK)
    img = IMG_DIR.parent / "feature_maps.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(5.8), Inches(3.0), Inches(3.8))

    # Phase 3 detail
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Phase 3: Analysis & Modeling")
    add_bullet_slide(slide, Inches(0.6), Inches(1.2), Inches(5.0), Inches(1.5), [
        "Explored distributions, correlations, and VIF across 27 features",
        "Found strong multicollinearity: restaurant ≈ evening businesses,\n"
        "SES cluster ≈ income, noise ≈ infrastructure complaints",
    ], font_size=13, spacing=Pt(4))
    add_text_box(slide, Inches(0.6), Inches(2.8), Inches(5.0), Inches(0.4),
                 "Three modeling approaches (5-fold CV):", font_size=14, bold=True, color=DARK)
    models = [
        ("Linear Reg.", "Stepwise, 17 features", "R² = 0.330"),
        ("Neural Net", "2 layers (64-32), early stop", "R² = 0.437"),
        ("Random Forest", "300 trees, leaf=5", "R² = 0.489"),
    ]
    for i, (name, desc, r2) in enumerate(models):
        y = 3.25 + i * 0.42
        clr = ACCENT if i == 2 else DARK
        add_text_box(slide, Inches(0.8), Inches(y), Inches(1.5), Inches(0.35),
                     name, font_size=12, bold=True, color=clr)
        add_text_box(slide, Inches(2.3), Inches(y), Inches(2.5), Inches(0.35),
                     desc, font_size=11, color=GRAY)
        add_text_box(slide, Inches(4.6), Inches(y), Inches(1.2), Inches(0.35),
                     r2, font_size=12, bold=True, color=clr)

    add_text_box(slide, Inches(0.6), Inches(4.6), Inches(5.0), Inches(0.6),
                 "PCA walkability composite score (0–100):\n"
                 "PC1 captures 26% of variance across 24 features",
                 font_size=12, color=DARK)

    img = IMG_DIR / "02_correlation.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(5.8), Inches(1.2), Inches(4.0))

    # Phase 4 detail
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Phase 4: Validation & Interpretation")
    add_bullet_slide(slide, Inches(0.6), Inches(1.2), Inches(5.0), Inches(4.0), [
        "Permutation feature importance: restaurant_count,\n"
        "cars_per_100_residents, and ses_cluster dominate",
        "Mapped all features to the paper's 6 VWP categories\n"
        "to compare our findings against visual-feature results",
        "PCA walkability composite score (0–100 scale):\n"
        "PC1 explains 26% of variance across 24 features",
        "Built interactive Folium map with click-to-inspect\n"
        "walkability scores for each street segment",
        "Generated policy recommendations for low-scoring areas",
    ], font_size=13, spacing=Pt(4))
    img = IMG_DIR / "07_permutation_importance.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(5.5), Inches(1.2), Inches(4.3))

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 8. RESULTS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 8, "Results")

    # Dataset overview
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "The Dataset: 1,278 Segments × 26 Features")
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(0.5),
                 "2,070 segments reduced to 1,278 after dropping rows with missing values",
                 font_size=14, color=GRAY)
    add_text_box(slide, Inches(0.6), Inches(1.8), Inches(8.8), Inches(0.4),
                 "What one row looks like — three example segments:",
                 font_size=14, bold=True, color=DARK)
    # Table header
    cols = [("Feature", 2.5), ("Rothschild Blvd\n(high walkability)", 2.2),
            ("Ben Tzvi\n(mid)", 1.8), ("Yerushalayim St\n(low walkability)", 2.2)]
    for j, (hdr, w) in enumerate(cols):
        x = 0.6 + sum(c[1] for c in cols[:j])
        add_text_box(slide, Inches(x), Inches(2.3), Inches(w), Inches(0.55),
                     hdr, font_size=10, bold=True, color=WHITE,
                     alignment=PP_ALIGN.CENTER)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(x), Inches(2.3), Inches(w), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        for k, line in enumerate(hdr.split("\n")):
            if k == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            _fmt_para(p, line, 10, True, WHITE, "Calibri", PP_ALIGN.CENTER)

    rows = [
        ("Neighborhood", "Lev Tel-Aviv", "Machon Yaffo TLV", "Yaffo G + Neve Golan"),
        ("avg_speed", "11.6 km/h", "13.0 km/h", "17.3 km/h"),
        ("restaurant_count", "635", "19", "1"),
        ("total_businesses", "4,806", "251", "54"),
        ("ses_cluster", "6", "2", "1"),
        ("cars_per_100_residents", "39", "37", "34"),
        ("noise_complaints", "5,270", "485", "760"),
        ("speed_variance", "1.3", "4.1", "0.1"),
    ]
    for i, (feat, v1, v2, v3) in enumerate(rows):
        yy = 2.85 + i * 0.3
        bg = LIGHT_BG if i % 2 == 0 else WHITE
        vals = [feat, v1, v2, v3]
        for j, (val, (_, w)) in enumerate(zip(vals, cols)):
            x = 0.6 + sum(c[1] for c in cols[:j])
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(x), Inches(yy), Inches(w), Inches(0.3))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.fill.background()
            is_label = (j == 0)
            tf = shape.text_frame
            tf.word_wrap = True
            _fmt_para(tf.paragraphs[0], val, 9, is_label,
                      DARK, "Calibri",
                      PP_ALIGN.LEFT if is_label else PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(5.2), Inches(8.8), Inches(0.3),
                 "Walkable streets: slow traffic, many restaurants, high SES.  "
                 "Non-walkable: faster traffic, few businesses, low SES.",
                 font_size=11, bold=True, color=ACCENT)

    # Key finding
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Municipal Data Captures Walkability")
    add_text_box(slide, Inches(0.6), Inches(1.3), Inches(8.8), Inches(0.8),
                 "Municipal features alone explain more traffic speed variance\n"
                 "than the traffic's own temporal patterns.",
                 font_size=18, color=DARK)
    add_text_box(slide, Inches(0.6), Inches(2.3), Inches(8.8), Inches(0.5),
                 "Combined model (traffic + municipal): Adj-R² = 0.354",
                 font_size=22, bold=True, color=ACCENT)
    add_text_box(slide, Inches(0.6), Inches(3.0), Inches(8.8), Inches(0.4),
                 "Top 5 most important features:", font_size=16, bold=True, color=DARK)
    top5 = [
        ("1.  Restaurant count", "— street-level vitality signal"),
        ("2.  Cars per 100 residents", "— car dependency = less walkable"),
        ("3.  SES cluster", "— wealthier areas have better infrastructure"),
        ("4.  Noise complaints", "— direct pedestrian discomfort"),
        ("5.  Total businesses", "— mixed-use streets attract walking"),
    ]
    for i, (feat, desc) in enumerate(top5):
        y = 3.4 + i * 0.38
        add_text_box(slide, Inches(1.0), Inches(y), Inches(3.5), Inches(0.35),
                     feat, font_size=14, bold=True, color=DARK)
        add_text_box(slide, Inches(4.5), Inches(y), Inches(5.0), Inches(0.35),
                     desc, font_size=14, color=GRAY)

    # Walkability map
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Tel Aviv Walkability Map")
    img = IMG_DIR / "05_walkability_map.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.3), Inches(1.1), Inches(5.5))
    add_text_box(slide, Inches(5.9), Inches(1.3), Inches(3.8), Inches(0.4),
                 "Most walkable", font_size=16, bold=True, color=ACCENT)
    add_bullet_slide(slide, Inches(5.9), Inches(1.7), Inches(3.8), Inches(1.2), [
        "Rothschild Blvd (100/100)",
        "Old North — south part (73)",
        "Old North — north part (62)",
        "Montefiori (48)",
    ], font_size=13, color=DARK, spacing=Pt(2))
    add_text_box(slide, Inches(5.9), Inches(3.1), Inches(3.8), Inches(0.4),
                 "Least walkable", font_size=16, bold=True, color=RED)
    add_bullet_slide(slide, Inches(5.9), Inches(3.5), Inches(3.8), Inches(1.2), [
        "Neve Dan (3/100)",
        "Rabiviim (5)",
        "Ramat HaChayal (9)",
        "Neot Afeka (10)",
    ], font_size=13, color=DARK, spacing=Pt(2))
    add_text_box(slide, Inches(5.9), Inches(4.7), Inches(3.8), Inches(0.5),
                 "Pattern: central mixed-use → high;\nperipheral residential → low",
                 font_size=12, color=GRAY)

    # Comparison with paper
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "How Our Findings Compare to the Paper")
    comparisons = [
        ("Comfort", "Paper: car/truck pixels (–)\nOurs: cars per 100 residents, noise complaints",
         "Strong match"),
        ("Pleasurability", "Paper: person/bicycle pixels (+)\nOurs: restaurants, events, young adults",
         "Strong match"),
        ("Safety", "Paper: vegetation (+), truck (–)\nOurs: female pop growth, evening businesses",
         "Novel signal"),
        ("Accessibility", "Paper: sidewalk (+), fence (–)\nOurs: infrastructure complaints, construction",
         "Good proxy"),
    ]
    for i, (cat, finding, match) in enumerate(comparisons):
        y = 1.2 + i * 1.0
        add_text_box(slide, Inches(0.6), Inches(y), Inches(1.5), Inches(0.4),
                     cat, font_size=14, bold=True, color=ACCENT)
        add_text_box(slide, Inches(2.2), Inches(y), Inches(5.5), Inches(0.8),
                     finding, font_size=12, color=DARK)
        match_color = ACCENT if "Strong" in match else BLUE
        add_text_box(slide, Inches(8.0), Inches(y), Inches(1.8), Inches(0.4),
                     match, font_size=12, bold=True, color=match_color)

    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(8.8), Inches(0.4),
                 "Novel: female population growth as safety proxy has no visual equivalent",
                 font_size=13, bold=True, color=DARK)

    # Neighborhood ranking
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Neighborhood Walkability Ranking")
    img = IMG_DIR / "06_neighborhood_ranking.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.3), Inches(1.1), Inches(5.2))
    add_text_box(slide, Inches(5.7), Inches(1.3), Inches(4.0), Inches(0.4),
                 "Key observations:", font_size=16, bold=True, color=DARK)
    add_bullet_slide(slide, Inches(5.7), Inches(1.8), Inches(4.1), Inches(3.5), [
        "Clear north-south SES gradient\nmatches walkability scores",
        "Central mixed-use neighborhoods\n(Lev TLV, Old North) score highest",
        "Car-dependent northern suburbs\n(Rabiviim, Neve Dan) score lowest",
        "Low-scoring areas share a pattern:\nfew restaurants, low events,\nhigh car ownership",
        "Jaffa neighborhoods are mid-range—\nhigh vitality but also high complaints",
    ], font_size=12, color=DARK, spacing=Pt(4))

    # Model comparison: three models
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Three Models Compared")
    img = IMG_DIR / "11_model_comparison.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.3), Inches(1.2), Inches(9.4), Inches(3.2))
    add_bullet_slide(slide, Inches(0.6), Inches(4.5), Inches(8.8), Inches(1.0), [
        "Random Forest (R²=0.489) > Neural Network (0.437) > Linear Regression (0.330)",
        "Neural net improves over linear (+32%) but is unstable with only 1,278 samples",
        "Random Forest wins: best accuracy, most stable, no tuning needed",
    ], font_size=12, color=DARK, spacing=Pt(2))

    # Feature importance comparison
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "Feature Importance Shifts Between Models")
    img = IMG_DIR / "12_feature_importance_comparison.png"
    if img.exists():
        slide.shapes.add_picture(str(img), Inches(0.2), Inches(1.2), Inches(9.6), Inches(3.3))
    add_bullet_slide(slide, Inches(0.6), Inches(4.6), Inches(8.8), Inches(0.8), [
        "Linear Regression: municipal features dominate (restaurant count, cars, SES)",
        "Random Forest: traffic features rise to top (speed variance, directional diff)",
        "Both agree: business count and speed patterns are key predictors",
    ], font_size=12, color=DARK, spacing=Pt(2))

    # Weighted Walkability Index formula
    slide = prs.slides.add_slide(blank)
    add_content_title(slide, "TLV Walkability Index — Weighted Formula")
    add_text_box(slide, Inches(0.6), Inches(1.2), Inches(9.0), Inches(0.5),
                 "Walkability  =  Σ ( wᵢ  ×  signᵢ  ×  z_scoreᵢ ),  normalized to 0–100",
                 font_size=14, bold=True, color=ACCENT)

    # Formula table
    formula_features = [
        ("restaurant_count",       "+", "0.31", "More restaurants = more street vitality"),
        ("cars_per_100_residents",  "−", "0.25", "Car dependency = less walkable"),
        ("ses_cluster",            "+", "0.21", "Higher SES = better infrastructure"),
        ("noise_complaints",       "−", "0.09", "Noise = pedestrian discomfort"),
        ("total_businesses",       "+", "0.05", "Mixed-use streets attract walking"),
        ("elderly_share",          "−", "0.03", "Higher elderly share = less active streets"),
        ("speed_variance",         "−", "0.02", "Unstable traffic = less safe"),
        ("dir_speed_diff",         "−", "0.02", "Asymmetric flow = bottleneck"),
        ("young_adult_share",      "+", "0.01", "Young adults prefer walkable areas"),
        ("female_pop_growth",      "+", "0.01", "Women move to safe neighborhoods"),
    ]
    # Header
    hdr_cols = [("Feature", 2.6), ("Sign", 0.6), ("Weight", 0.8), ("Interpretation", 4.8)]
    hdr_x = 0.5
    for hdr, w in hdr_cols:
        x = hdr_x + sum(c[1] for c in hdr_cols[:hdr_cols.index((hdr, w))])
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(x), Inches(1.75), Inches(w), Inches(0.32))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        _fmt_para(tf.paragraphs[0], hdr, 9, True, WHITE, "Calibri", PP_ALIGN.CENTER)
    # Rows
    for i, (feat, sign, weight, interp) in enumerate(formula_features):
        yy = 2.07 + i * 0.29
        bg = LIGHT_BG if i % 2 == 0 else WHITE
        sign_color = ACCENT if sign == "+" else RED
        vals = [feat, sign, weight, interp]
        for j, (val, (_, w)) in enumerate(zip(vals, hdr_cols)):
            x = hdr_x + sum(c[1] for c in hdr_cols[:j])
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(x), Inches(yy), Inches(w), Inches(0.29))
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            clr = sign_color if j == 1 else (ACCENT if j == 0 else DARK)
            _fmt_para(tf.paragraphs[0], val, 9, j == 0,
                      clr, "Calibri",
                      PP_ALIGN.CENTER if j in (1, 2) else PP_ALIGN.LEFT)

    add_text_box(slide, Inches(0.5), Inches(4.97), Inches(8.8), Inches(0.5),
                 "Weights derived from permutation importance of linear regression  |  "
                 "Correlation with PCA score: r = 0.950",
                 font_size=10, color=GRAY)

    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    # 9. APPENDIX — FEATURE DETAILS
    # ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
    slide = prs.slides.add_slide(blank)
    add_section_header(slide, 9, "Appendix: Feature Details")

    # --- Feature table helper ---
    def add_feature_table(slide, title, features, start_y=1.2):
        """Draw a compact feature table. features: list of (name, source, description)."""
        add_content_title(slide, title)
        # Header row
        headers = [("Feature", 2.8), ("Source / Join", 2.2), ("What It Captures", 4.6)]
        x_offset = 0.2
        for hdr, w in headers:
            x = x_offset + sum(c[1] for c in headers[:headers.index((hdr, w))])
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            Inches(x), Inches(start_y), Inches(w), Inches(0.35))
            shape.fill.solid()
            shape.fill.fore_color.rgb = DARK
            shape.line.fill.background()
            tf = shape.text_frame
            tf.word_wrap = True
            _fmt_para(tf.paragraphs[0], hdr, 10, True, WHITE, "Calibri", PP_ALIGN.CENTER)

        for i, (fname, source, desc) in enumerate(features):
            yy = start_y + 0.35 + i * 0.33
            bg = LIGHT_BG if i % 2 == 0 else WHITE
            vals = [fname, source, desc]
            for j, (val, (_, w)) in enumerate(zip(vals, headers)):
                x = x_offset + sum(c[1] for c in headers[:j])
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                Inches(x), Inches(yy), Inches(w), Inches(0.33))
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg
                shape.line.fill.background()
                tf = shape.text_frame
                tf.word_wrap = True
                is_label = (j == 0)
                _fmt_para(tf.paragraphs[0], val, 9, is_label,
                          ACCENT if is_label else DARK, "Calibri",
                          PP_ALIGN.LEFT if is_label else PP_ALIGN.LEFT)

    # --- Slide: Traffic + Business features (10) ---
    slide = prs.slides.add_slide(blank)
    add_feature_table(slide, "Traffic Features (5) + Business Features (5)", [
        ("avg_speed", "Hourly speed cols", "Mean of 15 hourly speeds (6am-8pm)"),
        ("peak_speed_drop", "Hourly speed cols", "Ratio: 8am speed / 6am speed (rush impact)"),
        ("speed_variance", "Hourly speed cols", "Variance across 15 hours (stability)"),
        ("dir_speed_diff", "Both directions", "|speed dir A - speed dir B| (asymmetry)"),
        ("evening_speed_recovery", "Hourly speed cols", "Ratio: 7pm speed / 5pm speed"),
        ("restaurant_count", "Businesses / nh_code", "Food & hospitality businesses per neighborhood"),
        ("total_businesses", "Businesses / nh_code", "All registered businesses per neighborhood"),
        ("biz_diversity_index", "Businesses / nh_code", "Shannon entropy of business categories"),
        ("licensed_biz_ratio", "Businesses / nh_code", "Fraction requiring a license (formality)"),
        ("evening_biz_count", "Businesses / nh_code", "Restaurants + entertainment (nightlife)"),
    ])

    # --- Slide: Street Events + Complaints + Construction (9) ---
    slide = prs.slides.add_slide(blank)
    add_feature_table(slide, "Street Events (5) + Complaints (2) + Construction (2)", [
        ("event_frequency", "Closed streets / street_code", "Number of event-type closures on this street"),
        ("event_diversity", "Closed streets / street_code", "Distinct event types on this street"),
        ("light_rail_closures", "Closed streets / street_code", "Light rail construction closures"),
        ("parade_route_count", "Closed streets / street_code", "Marathon / parade route closures"),
        ("infra_closure_count", "Closed streets / street_code", "Infrastructure work closures"),
        ("noise_complaints", "Hotline 106 / nh_code", "Noise-related complaints per neighborhood"),
        ("infra_complaints", "Hotline 106 / nh_code", "Sidewalk/road/lighting complaints"),
        ("construction_starts", "CBS / sub_district→nh", "Building starts (latest year)"),
        ("construction_completion_rate", "CBS / sub_district→nh", "Completions / starts ratio"),
    ])

    # --- Slide: Demographics + Housing + SES (8) ---
    slide = prs.slides.add_slide(blank)
    add_feature_table(slide, "Demographics (3) + Housing (2) + Socio-Economic (3)", [
        ("female_pop_growth_yoy", "Population / name match", "Year-over-year female population growth"),
        ("young_adult_share", "Population / name match", "% of population aged 20-34"),
        ("elderly_share", "Population / name match", "% of population aged 65+"),
        ("residential_density", "Dwelling / nh_code", "Dwelling units per billed area"),
        ("commercial_to_residential_ratio", "Dwelling / nh_code", "Commercial units / residential units"),
        ("ses_cluster", "SES / name + manual map", "Socio-economic cluster (1-10)"),
        ("cars_per_100_residents", "SES / name + manual map", "Car ownership rate"),
        ("avg_monthly_income_per_capita", "SES / name + manual map", "Average monthly income"),
    ])

    # ── CLOSING SLIDE ───────────────────────────────────────────────────────
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, DARK)
    add_text_box(slide, Inches(0.8), Inches(1.0), Inches(8.4), Inches(1.0),
                 "Key Takeaway",
                 font_size=36, bold=True, color=WHITE)
    add_text_box(slide, Inches(0.8), Inches(2.0), Inches(8.4), Inches(1.5),
                 "A city's own data — business registries, complaint logs,\n"
                 "demographics — reveals how walkable its streets feel,\n"
                 "capturing dimensions that even street-level photos miss.",
                 font_size=20, color=ACCENT)
    add_text_box(slide, Inches(0.8), Inches(3.8), Inches(8.4), Inches(1.0),
                 "No cameras needed. No surveys. Just open data.",
                 font_size=18, color=WHITE)
    add_text_box(slide, Inches(0.8), Inches(4.7), Inches(8.4), Inches(0.5),
                 "Thank you  |  Questions?", font_size=16, color=GRAY)

    # ── SAVE ────────────────────────────────────────────────────────────────
    prs.save(str(OUT_PATH))
    print(f"Presentation saved: {OUT_PATH}")
    print(f"  {len(prs.slides)} slides")


if __name__ == "__main__":
    build()
